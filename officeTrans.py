import json
import docx
from tencentcloud.common import credential
from tencentcloud.common.profile.client_profile import ClientProfile
from tencentcloud.common.profile.http_profile import HttpProfile
from tencentcloud.common.exception.tencent_cloud_sdk_exception import TencentCloudSDKException
from tencentcloud.tmt.v20180321 import tmt_client, models
import openpyxl
from openpyxl.styles import Alignment
import os
import time
import datetime
from tkinter import filedialog
import tkinter as tk
from tkinter import ttk
from pptx import Presentation

lag = 0.25  #此参数为了配合腾讯云免费账户的翻译限制，表示1秒钟调用4次API，如果是付费账户可以改为0
langDic = {
    '中文': 'zh',
    '英文': 'en',
    '日语': 'jp',
    '韩语': 'kr',
    '德语': 'de',
    '法语': 'fr',
    '西班牙文': 'es',
    '意大利文': 'it',
    '俄文': 'ru',
    '葡萄牙文': 'pt',
    '越南文': 'vi',
    '印度尼西亚文': 'id',
    '马来西亚文': 'ms',
    '泰文': 'th'
}
charnum = 0
wordnum = 0
cwordnum = 0

# 调用API进行翻译
def translateTencent(content, fromLang, toLang, appId, secretKey):
    print(content, fromLang, toLang, appId, secretKey)
    if(len(content) < 5 and content != 'we'):  # 小于5个字符的不翻译,we是为了测试API是否正常
        return content
    try:
        cred = credential.Credential(appId, secretKey)
        httpProfile = HttpProfile()
        httpProfile.endpoint = "tmt.tencentcloudapi.com"
        clientProfile = ClientProfile()
        clientProfile.httpProfile = httpProfile
        client = tmt_client.TmtClient(cred, "ap-shanghai", clientProfile)
        req = models.TextTranslateRequest()
        params = {
            "SourceText": content,
            "Source": fromLang,
            "Target": toLang,
            "DocumentType": 'docx',
            'ProjectId': 0,
            "UntranslateTencentdText": "RBA"
        }
        req.from_json_string(json.dumps(params))
        resp = client.TextTranslate(req).TargetText
        return resp
    except TencentCloudSDKException as err:
        print(err)

# Excel的翻译
def excelTrans(srcFilename, fromLang, toLang, isReserve, appId, secretKey):
    rex = '.xlsx'
    sname = srcFilename[:srcFilename.rfind(".")]
    wb = openpyxl.load_workbook(srcFilename)
    sheets = wb.sheetnames
    # 遍历每个Sheet进行翻译
    print("需要翻译的所有sheet列表")
    print(sheets)
    for i in range(len(sheets)):
        ws = wb[sheets[i]]
        print('正在翻译'+sheets[i]+',请耐心等待......')
        for r in range(1, ws.max_row + 1):
            for c in range(1, ws.max_column + 1):
                b = ws.cell(row=r, column=c).value
                if b is not None and len(str(b)) > 1 and "=" not in str(b):
                    global charnum
                    charnum += len(str(b))
                    if isReserve == 'yes':
                        result = str(b)+"\n" + \
                            translateTencent(str(b), fromLang,
                                             toLang, appId, secretKey)
                    else:
                        result = translateTencent(
                            str(b), fromLang, toLang, appId, secretKey)

                    ws.cell(row=r, column=c).value = str(result)

                    align = Alignment(horizontal='left',
                                      vertical='center', wrap_text=True)

                    ws.cell(row=r, column=c).alignment = align

                    time.sleep(lag)
    wb.save(sname+'_translateTencentd'+rex)

# Word翻译-翻译段落内的文字
def replace_text_in_paragraphs(paragraphs, fromLang, toLang, appId, secretKey, isReserve,isprogress):
    global cwordnum
    global charnum
    for paragraph in paragraphs:
        if isprogress:
            cwordnum += 1
        if paragraph.text:
            runs = paragraph.runs
            # if cwordnum reaches 1000, print the progress of cwordnum/wordnum with a progress bar
            # print(cwordnum,wordnum)
            if cwordnum % 10 == 0 and cwordnum>0:
                print('▇'*int(cwordnum/wordnum*100/5)+' ' +
                      str(int(cwordnum/wordnum*100))+'%')
            charnum += len(str(paragraph.text))
            oldtext = paragraph.text
            new_text = translateTencent(
                oldtext, fromLang, toLang, appId, secretKey)
            if oldtext != new_text:
                if isReserve == 'yes':
                    new_text = str(paragraph.text)+"\n" + new_text
                runs[0].text = new_text if new_text else oldtext
                for run in runs[1:]:
                    r = run._element
                    r.getparent().remove(r)
            time.sleep(lag)

# Word翻译-翻译表格内的文字
def replace_text_in_tables(tables, fromLang, toLang, appId, secretKey, isReserve):
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.paragraphs:
                    replace_text_in_paragraphs(
                        cell.paragraphs, fromLang, toLang, appId, secretKey, isReserve,False)

# Word翻译-翻译文本框内的文字
def replace_text_in_textboxes(children, fromLang, toLang, appId, secretKey, isReserve):
    count = 0
    for child in children:
        # 通过类型判断目录
        if child.tag.endswith('txbx'):
            for ci in child.iter():
                if ci.tag.endswith('main}r'):
                    count += 1
                    global charnum
                    charnum += len(str(ci.text))
                    if ci.text:
                        new_text = str(translateTencent(
                            ci.text, fromLang, toLang, appId, secretKey))
                        if isReserve == 'yes':
                            new_text = str(ci.text)+"\n" + new_text
                        ci.text = str(new_text)
                        time.sleep(lag)

# Word翻译-依次翻译表格，文本框和段落内的内容
def wordTrans(srcFilename, fromLang, toLang, isReserve, appId, secretKey):
    sname = srcFilename[:srcFilename.rfind(".")]
    # 载入文档
    doc = docx.Document(srcFilename)
    global wordnum
    wordnum += len(doc.paragraphs)
    print('正在处理Word内部的表格，请耐心等待...')
    replace_text_in_tables(doc.tables, fromLang, toLang,
                           appId, secretKey, isReserve)

    # 遍历并替换所有文本框中的文本
    print('正在处理Word内部的文本框，请耐心等待...')
    replace_text_in_textboxes(doc.element.body.iter(),
                              fromLang, toLang, appId, secretKey, isReserve)
    # 遍历并替换所有段落中的文本
    print('开始翻译，请查看进度条...')
    replace_text_in_paragraphs(
        doc.paragraphs, fromLang, toLang, appId, secretKey, isReserve,True)

    # 保存替换后的文档
    doc.save(sname+"_translated.docx")

# PPT翻译
def pptTrans(srcFilename, fromLang, toLang, isReserve, appId, secretKey):
    ppt = Presentation(srcFilename)
    totalcharnum = 0
    for slide in ppt.slides:
        # 遍历幻灯片中所有文本框
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # 遍历文本框中的所有段落
            for paragraph in shape.text_frame.paragraphs:
                # 替换段落中的文本内容
                for run in paragraph.runs:
                    # 获取原始文本内容
                    text = run.text
                    totalcharnum += len(str(text))
    print('所有文字数量：', totalcharnum)
    # 遍历所有幻灯片
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                # 遍历原始文本框的段落和文本，并组合为一个字符串进行翻译
                combined_text = ""
                for p in text_frame.paragraphs:
                    combined_text += p.text.rstrip() + "\n"
                if 'Confidential property of Pentair' in combined_text.strip():
                    continue
                # 使用外部翻译 API 进行翻译
                if len(text_frame.paragraphs)>0 and len(text_frame.paragraphs[0].runs) > 0:
                    original_font = text_frame.paragraphs[0].runs[0].font
                translated_text = str(translateTencent(
                    combined_text, fromLang, toLang, appId, secretKey))
                # 去掉translated_text头部和尾部的所有换行符
                # 将翻译后的文本替换原始文本框的内容
                text_frame.clear()  # 清空原始文本框的内容
                text_frame.paragraphs[0]  # 删除原始文本框的第一个段落
                for k, v in enumerate(translated_text.split('\n')):
                    if k == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    if isReserve == 'yes':
                        p.text = str(combined_text.split('\n')[k])+"-" + v
                    else:
                        p.text = v
         
                    if len(p.runs)>0:
                        run = p.runs[0]
                        run.font.name = original_font.name
                        run.font.size = original_font.size
                        run.font.bold = original_font.bold
                        run.font.italic = original_font.italic
                        run.font.underline = original_font.underline
                        try:
                            # print(original_font.color.rgb)
                            run.font.color.rgb= original_font.color.rgb
                        except:
                            pass
                        try:
                            if 'NOT_THEME_COLOR' in str(original_font.color.theme_color):
                                pass
                            else:
                                run.font.color.theme_color= original_font.color.theme_color
                        except Exception as e:
                            # print(e)
                            pass
                        try:
                            p.alignment = original_font.alignment
                        except:
                            pass
                global charnum
                charnum += len(str(combined_text))
                # print(charnum)
                if charnum % 10 == 0:
                    print('▇'*int(charnum/totalcharnum*100/5)+' ' +
                          str(int(charnum/totalcharnum*100))+'%')
    modified_ppt_path = os.path.splitext(srcFilename)[0] + '_translated.pptx'
    ppt.save(modified_ppt_path)


# 选择文件并判断文件类型，从而调用不同的翻译方法
def runtrains(isReserve, fromLang, toLang, appId, secretKey):
    if not appId or not secretKey:
        print('请填写secredId和secretKey后再运行')
        return
    if fromLang == toLang:
        print('源语言和目标语言不能相同，请重新选择')
        return
    testword = translateTencent('we', 'en', 'zh', appId, secretKey)
    # 先翻译一个简单的单词测试下是否调用API正常
    if testword != '我们':
        print('secredId和secretKey不正确或无法使用，请重新填写')
        return
    print('腾讯云连接成功，开始翻译...')
    start = datetime.datetime.now()
    my_filetypes = [('MS-Office files', '.xlsx'),
                    ('MS-Office files', '.xlsm'),
                    ('MS-Office files', '.docx'),
                    ('MS-Office files', '.docm'),
                    ('MS-Office files', '.pptx'),
                    ('MS-Office files', '.pptm')]
    answer = filedialog.askopenfilename(initialdir=os.getcwd(),
                                        title="只支持新版Office文件，如老版如xls的请先用excel转成xlsx再运行:",
                                        filetypes=my_filetypes)

    if answer:
        root.destroy()
        file_name = os.path.basename(answer)
        fileType = file_name.split('.')[-1].upper()
        print('翻译速度取决于表格复杂度和文字量，请耐心等待，提前关闭不能保存翻译内容，请自行评估时间...')
        if fileType == 'XLSM' or fileType == 'XLSX':
            excelTrans(answer, fromLang, toLang, isReserve, appId, secretKey)
        elif fileType == 'DOCX' or fileType == 'DOCM':
            wordTrans(answer, fromLang, toLang, isReserve, appId, secretKey)
        else:
            pptTrans(answer, fromLang, toLang, isReserve, appId, secretKey)
        print('翻译成功，请查看软件同目录下的翻译完成文件...')
        end = datetime.datetime.now()
        s = end-start
        global charnum
        global wordnum
        global cwordnum
        print('本次翻译共翻译'+str(charnum)+'个字符,总共花费时间'+str(s)+'秒')
        charnum = 0
        wordnum = 0
        cwordnum = 0
        os.system("pause")

def on_entry1_click(event):
    if entry_var1.get() == "请输入SecretId":
        entry1.delete(0, "end")


def on_entry2_click(event):
    if entry_var2.get() == "请输入SecretId":
        entry2.delete(0, "end")

 
if __name__ == '__main__':
    print('正在打开程序，请耐心等待，此窗口会显示一些运行信息，使用过程中不要关闭此窗口')
    root = tk.Tk()
    root.title("By Michael.Z")
    root.columnconfigure(0, minsize=100)
    root.columnconfigure(1, minsize=100)
    root.configure(padx=10, pady=10)
    # 第一行：标题
    title_label = tk.Label(root, text="Office文档翻译", font=("Arial", 28))
    title_label.grid(row=0, column=0, columnspan=20, pady=20)

    # 第二行：两个选项radio
    radio_var1 = tk.StringVar()
    radio_var1.set("no")
    radio_yes = tk.Radiobutton(
        root, text="保留原文", variable=radio_var1, value="yes")
    radio_no = tk.Radiobutton(
        root, text="不保留原文", variable=radio_var1, value="no")
    radio_yes.grid(row=2, column=0, padx=10, pady=5)
    radio_no.grid(row=2, column=1, padx=10, pady=5)

    # 第三行：两个combo并排在一起
    cb1_label = tk.Label(root, text="原文语言:")
    cb1_label.grid(row=3, column=0)
    combo1_var = tk.StringVar()
    combo1_var.set("英文")
    combo1 = ttk.Combobox(root, textvariable=combo1_var)
    combo1["values"] = list(langDic.keys())
    combo1.grid(row=3, column=1, padx=10, pady=5)
    combo1.configure(state="readonly")

    cb2_label = tk.Label(root, text="译文语言:")
    cb2_label.grid(row=4, column=0)
    combo2_var = tk.StringVar()
    combo2_var.set("中文")
    combo2 = ttk.Combobox(root, textvariable=combo2_var)
    combo2["values"] = list(langDic.keys())
    combo2.grid(row=4, column=1, padx=10, pady=5)
    combo2.configure(state="readonly")

    # 第四行：一个按钮
    button = tk.Button(root, text="选择文件开始翻译", command=lambda: runtrains(radio_var1.get(), langDic[combo1_var.get()],
                                                                            langDic[combo2_var.get()], entry_var1.get(), entry_var2.get()), bg="blue", fg="white")
    button.grid(row=6, column=0, columnspan=2, padx=10, pady=20)

    # 第五行：一个标签
    bottom_label = tk.Label(root, text="*本翻译基于腾讯云机器翻译\n网址：https://cloud.tencent.com/product/tmt,\n请自行注册,并在下方填入secretID和secretKey*")
    bottom_label.grid(row=7, column=0, columnspan=2, padx=10, pady=5)

    entry_var1 = tk.StringVar()
    entry1 = tk.Entry(root, textvariable=entry_var1)
    entry1.bind("<FocusIn>", on_entry1_click)
    entry1.insert(0, "请输入SecretId")
    entry1.grid(row=8, column=0, columnspan=5, pady=5, sticky="we")
    entry1.configure(fg="blue")

    entry_var2 = tk.StringVar()
    entry2 = tk.Entry(root, textvariable=entry_var2, show="*")
    entry2.bind("<FocusIn>", on_entry2_click)
    entry2.insert(0, "请输入SecretId")
    entry2.grid(row=9, column=0, columnspan=5, sticky="we")
    entry2.configure(fg="blue")

    screenWidth = root.winfo_screenwidth() # 获取显示区域的宽度
    screenHeight = root.winfo_screenheight() # 获取显示区域的高度
    width = 300 # 设定窗口宽度
    height = 400 # 设定窗口高度
    left = (screenWidth - width) / 2
    top = (screenHeight - height) / 2

    root.geometry("%dx%d+%d+%d" % (width, height, left, top))
    # 禁止root的缩放
    root.resizable(False, False)
    root.mainloop()
